from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Light-mode palette ──────────────────────────────────────────────────────
BG       = RGBColor(0xF8, 0xFA, 0xFC)  # slide background
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
NAVY     = RGBColor(0x0A, 0x16, 0x28)  # header bar, titles
CHARCOAL = RGBColor(0x1E, 0x29, 0x3B)  # body text
SLATE    = RGBColor(0x64, 0x74, 0x8B)  # secondary text
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)  # card borders / dividers
CARD     = RGBColor(0xFF, 0xFF, 0xFF)  # card fill
TEAL     = RGBColor(0x0D, 0x94, 0x88)  # built accent
TEAL_BG  = RGBColor(0xCC, 0xEB, 0xE9)  # teal badge fill
GREEN    = RGBColor(0x15, 0x80, 0x3D)
AMBER    = RGBColor(0xB4, 0x5A, 0x09)  # darker amber for light bg
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
RED      = RGBColor(0xB9, 0x1C, 0x1C)  # darker red for light bg
RED_BG   = RGBColor(0xFE, 0xE2, 0xE2)
PILL_BG  = RGBColor(0xF1, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line_color=None, line_w=0.75):
    s = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=12, bold=False, italic=False,
        color=CHARCOAL, align=PP_ALIGN.LEFT, wrap=True):
    tf = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def header_bar(slide, title, subtitle):
    """Navy top bar with white text."""
    box(slide, 0, 0, 13.33, 1.05, fill=NAVY)
    box(slide, 0, 0, 0.07, 1.05, fill=TEAL)          # teal left accent
    txt(slide, title,    0.22, 0.12, 10, 0.46,
        size=24, bold=True, color=WHITE)
    txt(slide, subtitle, 0.22, 0.60, 10, 0.30,
        size=11, color=RGBColor(0xA0, 0xAE, 0xC0))


def footer(slide):
    box(slide, 0, 7.28, 13.33, 0.22, fill=NAVY)
    txt(slide, "GGHN STARR  ·  starr@gghn.org.ng  ·  April 2026",
        0.22, 7.30, 9, 0.18, size=8,
        color=RGBColor(0xA0, 0xAE, 0xC0))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  What's Built vs Pending
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, BG)
header_bar(s1,
    "GGHN STARR — Platform Status",
    "Current Build Analysis  ·  April 2026")

# ── stat pills ───────────────────────────────────────────────────────────────
pills = [
    ("~35%", "Built to Date",    TEAL,  TEAL_BG),
    ("6",    "Pages Live",       GREEN, RGBColor(0xDC, 0xFC, 0xE7)),
    ("3",    "Partial Features", AMBER, AMBER_BG),
    ("9",    "Pending Features", RED,   RED_BG),
]
for i, (num, label, col, pbg) in enumerate(pills):
    x = 0.22 + i * 3.26
    box(s1, x, 1.14, 3.1, 0.72, fill=WHITE,
        line_color=BORDER, line_w=0.75)
    box(s1, x, 1.14, 3.1, 0.06, fill=col)   # top color stripe
    txt(s1, num,   x+0.14, 1.22, 2.7, 0.34,
        size=22, bold=True, color=col)
    txt(s1, label, x+0.14, 1.55, 2.7, 0.24,
        size=9, color=SLATE)

# ── section label row ────────────────────────────────────────────────────────
box(s1, 0.22, 1.96, 6.2, 0.26, fill=RGBColor(0xE6, 0xF7, 0xF5))   # teal tint
txt(s1, "BUILT", 0.36, 1.98, 5.8, 0.22,
    size=8, bold=True, color=TEAL)

box(s1, 6.6, 1.96, 6.5, 0.26, fill=RED_BG)
txt(s1, "PARTIAL / PENDING", 6.74, 1.98, 6.2, 0.22,
    size=8, bold=True, color=RED)

# vertical divider
box(s1, 6.47, 1.96, 0.02, 5.3, fill=BORDER)

# ── LEFT: built items ────────────────────────────────────────────────────────
built = [
    ("Homepage",
     "Hero · rotating stats · 3 pillars · partner strip · newsletter signup"),
    ("Policy Briefs  /briefs",
     "Filterable grid (month & theme) — 3 briefs live with PDF & infographic downloads"),
    ("Brief Detail pages",
     "Executive summary · key messages · author bio · prev-next navigation"),
    ("Experts  /experts",
     "3 expert profile cards with specialties and bios"),
    ("Methodology  /methodology",
     "Tabbed explainer: SEIR models · NIPAD platform · GlobalPPS & WHONET"),
    ("Contact  /contact",
     "Formspree form with inquiry-type dropdown"),
    ("Infrastructure",
     "Static export · Tailwind v4 design system · GAS newsletter + Formspree"),
]

for i, (title, detail) in enumerate(built):
    y = 2.30 + i * 0.68
    # teal dot
    box(s1, 0.26, y+0.07, 0.12, 0.12, fill=TEAL)
    txt(s1, title,  0.46, y,      5.8, 0.24,
        size=10, bold=True, color=CHARCOAL)
    txt(s1, detail, 0.46, y+0.24, 5.8, 0.34,
        size=8.5, color=SLATE)

# ── RIGHT: pending items ─────────────────────────────────────────────────────
rows = [
    (AMBER, AMBER_BG, "PARTIAL", "Awareness Hub",
     "Briefs exist; no infographic gallery or explainer articles"),
    (AMBER, AMBER_BG, "PARTIAL", "Education Library",
     "No audience segmentation (Public / HCW / Policymakers)"),
    (AMBER, AMBER_BG, "PARTIAL", "Social Media",
     "LinkedIn in footer only — no share buttons on content"),
    (RED, RED_BG, "MISSING", "Audience-Segmented CTAs",
     "Single generic hero; no separate entry paths per audience"),
    (RED, RED_BG, "MISSING", "News Section",
     "/news listing + article pages not yet started"),
    (RED, RED_BG, "MISSING", "Take Action Page",
     "No pledges, prescribing commitments, or advocacy toolkit"),
    (RED, RED_BG, "MISSING", "Interactive AMR Data Map",
     "Mapping library + country-level AMR data not integrated"),
    (RED, RED_BG, "MISSING", "Practical Tools",
     "Stewardship checklists · quizzes · facility templates"),
    (RED, RED_BG, "MISSING", "Analytics",
     "No GA4, Plausible, or event tracking wired in"),
    (RED, RED_BG, "MISSING", "Accessibility (a11y)",
     "No ARIA audit, skip-links, or screen-reader testing"),
]

for i, (col, cbg, status, title, detail) in enumerate(rows):
    y = 2.30 + i * 0.50
    # badge
    box(s1, 6.62, y+0.05, 0.88, 0.17, fill=cbg)
    tf2 = s1.shapes.add_textbox(
        Inches(6.63), Inches(y+0.06), Inches(0.87), Inches(0.15))
    p2 = tf2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = status
    r2.font.size = Pt(7)
    r2.font.bold = True
    r2.font.color.rgb = col

    txt(s1, title,  7.60, y,      5.5, 0.22,
        size=10, bold=True, color=CHARCOAL)
    txt(s1, detail, 7.60, y+0.22, 5.5, 0.22,
        size=8, color=SLATE)

footer(s1)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  Priority Roadmap
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, BG)
header_bar(s2,
    "Recommended Build Priority",
    "Sequenced by platform impact  ·  High Priority first, then Medium")

# column headers
box(s2, 0.22, 1.14, 6.25, 0.28, fill=RED_BG)
txt(s2, "HIGH PRIORITY — Build Next Sprint", 0.36, 1.16, 6.0, 0.24,
    size=9, bold=True, color=RED)

box(s2, 6.72, 1.14, 6.38, 0.28, fill=AMBER_BG)
txt(s2, "MEDIUM PRIORITY — Following Sprints", 6.86, 1.16, 6.1, 0.24,
    size=9, bold=True, color=AMBER)

box(s2, 6.59, 1.10, 0.02, 6.2, fill=BORDER)


def cards(slide, items, x, border_col, badge_bg, badge_col):
    start_y = 1.50
    card_h  = 1.40
    gap     = 0.08

    for i, (title, bullets) in enumerate(items):
        y = start_y + i * (card_h + gap)
        box(slide, x, y, 6.25, card_h,
            fill=WHITE, line_color=BORDER, line_w=0.75)
        box(slide, x, y, 0.06, card_h, fill=border_col)  # left accent
        txt(slide, title, x+0.20, y+0.10, 5.9, 0.26,
            size=11, bold=True, color=CHARCOAL)
        for j, b in enumerate(bullets):
            txt(slide, "->  " + b,
                x+0.20, y+0.40+j*0.30, 5.9, 0.26,
                size=8.5, color=SLATE)


high = [
    ("Audience-Segmented CTAs", [
        "Rework hero with 3 distinct audience entry points",
        "Separate landing sections per audience (Public / HCW / Policy)",
        "Tag all existing content by target audience",
    ]),
    ("Analytics Integration", [
        "Add GA4 or Plausible to root layout",
        "Event tracking on downloads, CTAs, and newsletter signups",
        "Conversion funnel reporting from launch day",
    ]),
    ("News Section  (/news)", [
        "News listing page with category filters",
        "Article detail pages with publish date & author",
        "Social share buttons (LinkedIn, Twitter/X, WhatsApp)",
    ]),
    ("Take Action Page", [
        "Public pledge / commitment form",
        "Prescribing commitment tool for healthcare workers",
        "Downloadable advocacy toolkit for policymakers",
    ]),
]

med = [
    ("Awareness Hub & Education Library", [
        "Infographic gallery with filter by theme",
        "Short explainer articles per audience segment",
        "Audience filter across all resource pages",
    ]),
    ("Practical Tools Suite", [
        "Interactive stewardship checklist builder",
        "Self-assessment quiz flows with scored results",
        "Downloadable facility reporting templates",
    ]),
    ("Interactive AMR Data Map", [
        "Leaflet or Mapbox library integration",
        "Country-level AMR burden data overlay",
        "Drill-down to sub-national data where available",
    ]),
    ("Accessibility & Social Integration", [
        "Full ARIA landmark audit and skip-link",
        "Screen reader and colour-contrast testing",
        "Share buttons on all brief and news pages",
    ]),
]

cards(s2, high, 0.22, RED,   RED_BG,   RED)
cards(s2, med,  6.72, AMBER, AMBER_BG, AMBER)

footer(s2)


prs.save(r"c:\Users\HP_PC\Documents\gghn\gghnstarr\GGHN_STARR_Status.pptx")
print("Done.")
